from __future__ import annotations

import argparse
import hashlib
import json
import math
import re
from collections import Counter
from pathlib import Path

import numpy as np
import pandas as pd
from docx import Document
from openpyxl import load_workbook
from openpyxl.formatting.rule import FormulaRule
from openpyxl.styles import Alignment, Font, PatternFill
from openpyxl.utils import get_column_letter
from openpyxl.worksheet.datavalidation import DataValidation
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt


ROOT = Path(__file__).resolve().parents[1]
DEFAULT_MASTER = (
    ROOT
    / "artifacts"
    / "literature_agent_pce_stability"
    / "source_data"
    / "structured_extraction_validation_v1"
    / "structured_extraction_gold_standard_annotation_v1.xlsx"
)
DEFAULT_OUT = (
    ROOT
    / "artifacts"
    / "literature_agent_pce_stability"
    / "validation"
    / "structured_extraction_benchmark_v1"
)

YES_NO = {"yes", "no"}
YES_NO_PARTIAL = {"yes", "no", "partial"}
AUDIT_STATUSES = {"accepted", "sparse_accepted", "weak_accepted", "rejected"}

NUMERIC_TOLERANCES = [
    ("Cell_area_measured", 0.01, 0.05, "cm2", "max(0.01 cm2, 5% relative)"),
    ("JV_default_PCE", 0.10, 0.01, "%", "max(0.10 percentage point, 1% relative)"),
    ("JV_default_Voc", 0.01, 0.01, "V", "max(0.01 V, 1% relative)"),
    ("JV_default_Jsc", 0.20, 0.01, "mA cm-2", "max(0.20 mA cm-2, 1% relative)"),
    ("JV_default_FF", 0.01, 0.02, "fraction", "max(0.01 fraction, 2% relative); normalize percent to fraction"),
    ("Perovskite_deposition_thermal_annealing_temperature", 2.0, 0.02, "degC", "max(2 degC, 2% relative)"),
    ("Perovskite_deposition_thermal_annealing_time", 1.0, 0.05, "min", "max(1 min, 5% relative)"),
    ("Stability_time_total_exposure", 1.0, 0.02, "h", "max(1 h, 2% relative)"),
    ("Stability_temperature_range", 2.0, 0.02, "degC", "max(2 degC, 2% relative)"),
    ("Stability_relative_humidity_average_value", 2.0, 0.05, "% RH", "max(2 percentage points, 5% relative)"),
    ("Stability_light_intensity", 0.05, 0.05, "sun", "max(0.05 sun, 5% relative)"),
    ("Stability_PCE_initial_value", 0.10, 0.01, "%", "max(0.10 percentage point, 1% relative)"),
    ("Stability_PCE_end_of_experiment", 0.10, 0.02, "% or fraction", "max(0.10 point, 2% relative) after scale normalization"),
    ("Stability_PCE_T80", 5.0, 0.05, "h", "max(5 h, 5% relative)"),
    ("Stability_PCE_T95", 5.0, 0.05, "h", "max(5 h, 5% relative)"),
]

CATEGORICAL_FIELDS = {
    "Cell_architecture",
    "Stability_protocol",
}


def clean(value: object) -> str:
    if value is None or (isinstance(value, float) and math.isnan(value)):
        return ""
    text = str(value).strip()
    return "" if text.lower() in {"nan", "none", "null"} else text


def label(value: object) -> str:
    return clean(value).lower().replace(" ", "_")


def normalized_text(value: object) -> str:
    text = clean(value).lower()
    text = text.replace("\u2013", "-").replace("\u2014", "-").replace("\u2212", "-")
    return re.sub(r"\s+", " ", text).strip()


def canonical_category(field: str, value: object) -> str:
    text = normalized_text(value)
    compact = re.sub(r"[^a-z0-9]+", "", text)
    if field == "Cell_architecture":
        if "mesoporous" in compact:
            return "mesoporous"
        if "inverted" in compact or compact == "pin":
            return "p-i-n"
        if "conventional" in compact or "regular" in compact or compact == "nip":
            return "n-i-p"
    if field == "Stability_protocol":
        match = re.search(r"isos\s*[-_ ]?\s*([a-z])\s*[-_ ]?\s*(\d+)", text)
        if match:
            return f"ISOS-{match.group(1).upper()}-{match.group(2)}"
    return text or "__missing__"


def parse_number(value: object) -> float | None:
    text = normalized_text(value).replace(",", "")
    match = re.search(r"[-+]?\d*\.?\d+(?:e[-+]?\d+)?", text)
    if not match:
        return None
    try:
        return float(match.group(0))
    except ValueError:
        return None


def sha256(path: Path) -> str:
    digest = hashlib.sha256()
    with path.open("rb") as handle:
        for block in iter(lambda: handle.read(1024 * 1024), b""):
            digest.update(block)
    return digest.hexdigest()


def read_sheet(path: Path, name: str) -> pd.DataFrame:
    return pd.read_excel(path, sheet_name=name, dtype=object).fillna("")


def reorder_by_record(frame: pd.DataFrame, order: list[str]) -> pd.DataFrame:
    rank = {record_id: index for index, record_id in enumerate(order)}
    result = frame.copy()
    result["__record_order"] = result["validation_record_id"].map(rank)
    secondary = [column for column in ("field_family", "field") if column in result.columns]
    result = result.sort_values(["__record_order", *secondary], kind="stable")
    return result.drop(columns="__record_order").reset_index(drop=True)


def instruction_table() -> pd.DataFrame:
    instructions = [
        "Work independently and do not consult the other reviewer or the system integration status.",
        "Open the DOI landing page and use the full paper plus supplementary information when available.",
        "Confirm DOI and sample/device identity before scoring fields.",
        "Score every core field, including blank extracted values, so omissions are measured.",
        "Use only yes/no for source presence and yes/no/partial for correctness or evidence support.",
        "Use partial only when normalization preserves the scientific meaning but is not an exact representation.",
        "Enter a normalized gold value whenever present_in_source is yes; otherwise leave gold_value blank.",
        "For a blank extraction, mark correctness yes only when the field is genuinely absent from the source.",
        "For populated extractions, evidence support is yes only when the cited excerpt or inspected paper directly supports the value.",
        "Apply the frozen numeric tolerances; do not loosen them case by case.",
        "Review all additional populated fields for unsupported values, wrong units, or wrong sample assignment.",
        "Mark complete-row correctness yes only when identity and every assessed field in that row are correct.",
        "Assign accepted, sparse_accepted, weak_accepted, or rejected using the frozen decision rules.",
        "Use notes for ambiguity; do not resolve disagreements until both blinded files are complete.",
    ]
    return pd.DataFrame({"step": range(1, len(instructions) + 1), "instruction": instructions})


def tolerance_table() -> pd.DataFrame:
    return pd.DataFrame(
        NUMERIC_TOLERANCES,
        columns=["field", "absolute_tolerance", "relative_tolerance", "canonical_unit", "decision_rule"],
    )


def write_tables(path: Path, tables: dict[str, pd.DataFrame], input_prefixes: tuple[str, ...]) -> None:
    path.parent.mkdir(parents=True, exist_ok=True)
    with pd.ExcelWriter(path, engine="openpyxl") as writer:
        for sheet_name, frame in tables.items():
            frame.to_excel(writer, sheet_name=sheet_name[:31], index=False)
    workbook = load_workbook(path)
    header_fill = PatternFill("solid", fgColor="20242A")
    header_font = Font(name="Aptos", size=10, bold=True, color="FFFFFF")
    reference_fill = PatternFill("solid", fgColor="E7E8EA")
    input_fill = PatternFill("solid", fgColor="FFF2CC")
    r1_fill = PatternFill("solid", fgColor="DDEBF7")
    r2_fill = PatternFill("solid", fgColor="E2F0D9")
    adjudication_fill = PatternFill("solid", fgColor="FFF2CC")
    for worksheet in workbook.worksheets:
        worksheet.freeze_panes = "A2"
        worksheet.auto_filter.ref = worksheet.dimensions
        worksheet.sheet_view.showGridLines = False
        for cell in worksheet[1]:
            cell.fill = header_fill
            cell.font = header_font
            cell.alignment = Alignment(wrap_text=True, vertical="center")
        worksheet.row_dimensions[1].height = 34
        for column_index, column_cells in enumerate(worksheet.iter_cols(), start=1):
            header = clean(column_cells[0].value)
            values = [clean(cell.value) for cell in column_cells[:80]]
            width = min(max([len(header), *(len(value) for value in values), 10]) + 2, 48)
            if any(token in header for token in ("excerpt", "notes", "instruction", "decision_rule")):
                width = 45
            worksheet.column_dimensions[get_column_letter(column_index)].width = width
            if any(header.startswith(prefix) for prefix in input_prefixes):
                fill = adjudication_fill if header.startswith("adjudicated_") else input_fill
            elif header.startswith("r1_"):
                fill = r1_fill
            elif header.startswith("r2_"):
                fill = r2_fill
            else:
                fill = reference_fill
            for cell in column_cells[1:]:
                cell.fill = fill
                cell.font = Font(name="Aptos", size=9)
                cell.alignment = Alignment(wrap_text=True, vertical="top")
        headers = {clean(cell.value): cell.column for cell in worksheet[1]}
        for header, column_index in headers.items():
            letter = get_column_letter(column_index)
            if header.endswith("present_in_source_yes_no") or header.endswith("source_recovered_yes_no"):
                validation = DataValidation(type="list", formula1='"yes,no"', allow_blank=True)
                worksheet.add_data_validation(validation)
                validation.add(f"{letter}2:{letter}{worksheet.max_row}")
            elif any(
                header.endswith(suffix)
                for suffix in (
                    "correct_yes_no_partial",
                    "supports_value_yes_no_partial",
                    "supported_yes_no_partial",
                )
            ):
                validation = DataValidation(type="list", formula1='"yes,no,partial"', allow_blank=True)
                worksheet.add_data_validation(validation)
                validation.add(f"{letter}2:{letter}{worksheet.max_row}")
            elif any(
                header.endswith(suffix)
                for suffix in (
                    "sample_identity_correct_yes_no",
                    "sample_count_correct_yes_no",
                    "complete_core_row_correct_yes_no",
                    "complete_populated_row_correct_yes_no",
                )
            ):
                validation = DataValidation(type="list", formula1='"yes,no"', allow_blank=True)
                worksheet.add_data_validation(validation)
                validation.add(f"{letter}2:{letter}{worksheet.max_row}")
            elif header.endswith("recommended_audit_status"):
                validation = DataValidation(
                    type="list",
                    formula1='"accepted,sparse_accepted,weak_accepted,rejected"',
                    allow_blank=True,
                )
                worksheet.add_data_validation(validation)
                validation.add(f"{letter}2:{letter}{worksheet.max_row}")
        if "_source_landing_page" in headers:
            column = get_column_letter(headers["_source_landing_page"])
            for row in range(2, worksheet.max_row + 1):
                cell = worksheet[f"{column}{row}"]
                if clean(cell.value).startswith("http"):
                    cell.hyperlink = clean(cell.value)
                    cell.style = "Hyperlink"
        if "Ref_DOI_number" in headers:
            column = get_column_letter(headers["Ref_DOI_number"])
            for row in range(2, worksheet.max_row + 1):
                cell = worksheet[f"{column}{row}"]
                if clean(cell.value):
                    cell.hyperlink = f"https://doi.org/{clean(cell.value)}"
                    cell.style = "Hyperlink"
        if worksheet.max_row > 1 and input_prefixes:
            for prefix in input_prefixes:
                for header, column_index in headers.items():
                    if header.startswith(prefix):
                        letter = get_column_letter(column_index)
                        worksheet.conditional_formatting.add(
                            f"{letter}2:{letter}{worksheet.max_row}",
                            FormulaRule(
                                formula=[f'AND({letter}2<>"",COUNTIF({letter}2,"yes")+COUNTIF({letter}2,"no")+COUNTIF({letter}2,"partial")=0)'],
                                fill=PatternFill("solid", fgColor="F4CCCC"),
                            ),
                        )
    workbook.save(path)


def build_blinded(master: Path, out_dir: Path) -> dict[str, object]:
    manifest = read_sheet(master, "sample_manifest")
    core = read_sheet(master, "core_field_review")
    populated = read_sheet(master, "populated_field_review")
    record = read_sheet(master, "record_review")
    evidence = read_sheet(master, "evidence_lookup")
    rules = read_sheet(master, "decision_rules")

    safe_manifest_columns = [
        "validation_record_id",
        "Ref_DOI_number",
        "Ref_internal_sample_id",
        "Ref_lead_author",
        "Ref_publication_date",
        "Ref_journal",
        "_source_pdf_url",
        "_source_landing_page",
        "_full_text_source",
    ]
    safe_manifest = manifest[safe_manifest_columns].copy()
    base_core = [
        "validation_record_id",
        "Ref_DOI_number",
        "Ref_internal_sample_id",
        "field_family",
        "field",
        "extracted_value",
        "saved_evidence_excerpt",
    ]
    base_populated = [
        "validation_record_id",
        "Ref_DOI_number",
        "Ref_internal_sample_id",
        "field",
        "extracted_value",
        "saved_evidence_excerpt",
    ]
    base_record = ["validation_record_id", "Ref_DOI_number", "Ref_internal_sample_id"]

    reviewer_core = core[base_core].copy()
    for column in (
        "present_in_source_yes_no",
        "gold_value",
        "extracted_value_correct_yes_no_partial",
        "evidence_supports_value_yes_no_partial",
        "notes",
    ):
        reviewer_core[column] = ""
    reviewer_populated = populated[base_populated].copy()
    for column in ("correct_yes_no_partial", "supported_yes_no_partial", "notes"):
        reviewer_populated[column] = ""
    reviewer_record = record[base_record].copy()
    for column in (
        "source_recovered_yes_no",
        "sample_identity_correct_yes_no",
        "sample_count_correct_yes_no",
        "complete_core_row_correct_yes_no",
        "complete_populated_row_correct_yes_no",
        "recommended_audit_status",
        "notes",
    ):
        reviewer_record[column] = ""

    files: dict[str, str] = {}
    seeds = {"reviewer_1": 1401, "reviewer_2": 2401}
    record_ids = safe_manifest["validation_record_id"].astype(str).tolist()
    for reviewer, seed in seeds.items():
        rng = np.random.default_rng(seed)
        order = record_ids.copy()
        rng.shuffle(order)
        tables = {
            "instructions": instruction_table(),
            "sample_manifest": reorder_by_record(safe_manifest, order),
            "core_field_review": reorder_by_record(reviewer_core, order),
            "populated_field_review": reorder_by_record(reviewer_populated, order),
            "record_review": reorder_by_record(reviewer_record, order),
            "evidence_lookup": reorder_by_record(
                evidence[
                    [
                        "validation_record_id",
                        "Ref_DOI_number",
                        "Ref_internal_sample_id",
                        "field",
                        "evidence",
                        "lit_row_index",
                    ]
                ],
                order,
            ),
            "numeric_tolerances": tolerance_table(),
            "decision_rules": rules[["status", "operational_rule"]],
        }
        path = out_dir / f"structured_extraction_{reviewer}_blinded.xlsx"
        write_tables(path, tables, input_prefixes=("present_", "gold_", "extracted_", "evidence_", "correct_", "supported_", "source_", "sample_", "complete_", "recommended_", "notes"))
        files[reviewer] = str(path)

    protocol = out_dir / "ANNOTATION_PROTOCOL.md"
    protocol.write_text(
        """# Blinded structured-extraction validation protocol

## Frozen design

- 60 records from 59 DOI groups.
- 20 accepted PCE, 10 accepted stability, 15 sparse, 10 rejected, and 5 general accepted records in the hidden system key.
- At most two sampled records per DOI.
- 1,800 core-field decisions and 906 additional populated-field decisions.
- Both reviewers assess every record independently; the system status and sampling stratum are omitted from reviewer files.

## Review order

1. Confirm access to the paper and supplementary information.
2. Confirm DOI, sample identity, and the number of distinct device/sample records.
3. Complete every row in `core_field_review`, including blank extractions.
4. Complete every row in `populated_field_review`.
5. Complete `record_review` and assign an independent audit status.
6. Return the workbook without changing IDs, field names, sheet names, or column names.
7. Merge both files with the workflow script; an adjudicator then resolves disagreements without seeing the hidden system key.

## Primary scoring rule

The primary field metric is conservative and binary. A true positive requires a gold-present field, a populated extraction, and an adjudicated `yes` correctness decision. An incorrect populated value contributes both a false positive and a false negative. `partial` is reported separately and counts as incorrect in the primary metric. Empty fields that are truly absent are true negatives and do not inflate precision or recall.

## Evidence rule

Evidence support is scored only for populated values. `yes` requires direct support for the value and the same sample/device context. A topically related excerpt is not sufficient. `partial` is reserved for support that preserves the main meaning but omits a qualifier, unit, or sample discriminator.

## Numeric and categorical rules

Numeric comparisons use the frozen `numeric_tolerances` sheet after canonical-unit conversion. Categorical values use normalized exact matching with documented synonym normalization. The adjudicated reviewer decision remains authoritative when a source uses a representation not handled by automated parsing.

## Blinding

Reviewers must not inspect the master workbook, integration audit, or one another's decisions. The adjudicator sees both reviews but not the hidden system status. System audit agreement is calculated only after gold labels are frozen.
""",
        encoding="utf-8",
    )
    metric_spec = {
        "version": "1.0",
        "primary_partial_policy": "partial_counts_as_incorrect",
        "secondary_partial_policy": "partial_credit_reported_separately",
        "cluster_unit": "Ref_DOI_number",
        "bootstrap_iterations": 2000,
        "bootstrap_seed": 42,
        "categorical_fields": sorted(CATEGORICAL_FIELDS),
        "numeric_tolerances": [
            {
                "field": field,
                "absolute_tolerance": absolute,
                "relative_tolerance": relative,
                "canonical_unit": unit,
                "rule": rule,
            }
            for field, absolute, relative, unit, rule in NUMERIC_TOLERANCES
        ],
    }
    (out_dir / "metric_specification.json").write_text(
        json.dumps(metric_spec, indent=2), encoding="utf-8"
    )
    (out_dir / "RESULT_INSERTION_TEMPLATE.md").write_text(
        """# Structured-extraction result insertion template

The scoring command creates `manuscript_results_snippet.md`, an editable PowerPoint results figure, CSV/JSON metrics, and an updated manuscript only after every required adjudicated label is complete.

Primary manuscript outcomes:

- field value precision, recall, and F1 with DOI-cluster bootstrap intervals;
- evidence-support validity and unsupported-value rate;
- numeric tolerance agreement and categorical macro-F1;
- sample-linking and complete-row accuracy;
- audit-decision agreement;
- reviewer percent agreement and Cohen kappa.

Do not replace the manuscript's pending-result marker manually or report accuracy from incomplete reviewer files.
""",
        encoding="utf-8",
    )
    output_manifest = {
        "status": "awaiting_two_blinded_reviews",
        "master_workbook": str(master),
        "master_sha256": sha256(master),
        "reviewer_files": files,
        "reviewer_order_seeds": seeds,
        "records": int(len(safe_manifest)),
        "core_field_decisions_per_reviewer": int(len(reviewer_core)),
        "populated_field_decisions_per_reviewer": int(len(reviewer_populated)),
        "blinding": "integration status and validation stratum removed",
    }
    (out_dir / "blinded_validation_manifest.json").write_text(
        json.dumps(output_manifest, indent=2), encoding="utf-8"
    )
    return output_manifest


def validate_reviewer_values(frame: pd.DataFrame, columns: dict[str, set[str]], sheet: str) -> list[str]:
    errors: list[str] = []
    for column, allowed in columns.items():
        if column not in frame.columns:
            errors.append(f"{sheet}: missing column {column}")
            continue
        observed = {label(value) for value in frame[column] if clean(value)}
        invalid = sorted(observed - allowed)
        if invalid:
            errors.append(f"{sheet}.{column}: invalid values {invalid}")
    return errors


def merge_reviews(reviewer_1: Path, reviewer_2: Path, out_path: Path) -> dict[str, object]:
    sheets = ["core_field_review", "populated_field_review", "record_review"]
    r1 = {sheet: read_sheet(reviewer_1, sheet) for sheet in sheets}
    r2 = {sheet: read_sheet(reviewer_2, sheet) for sheet in sheets}
    errors: list[str] = []
    errors += validate_reviewer_values(
        r1["core_field_review"],
        {
            "present_in_source_yes_no": YES_NO,
            "extracted_value_correct_yes_no_partial": YES_NO_PARTIAL,
            "evidence_supports_value_yes_no_partial": YES_NO_PARTIAL,
        },
        "reviewer_1.core_field_review",
    )
    errors += validate_reviewer_values(
        r2["core_field_review"],
        {
            "present_in_source_yes_no": YES_NO,
            "extracted_value_correct_yes_no_partial": YES_NO_PARTIAL,
            "evidence_supports_value_yes_no_partial": YES_NO_PARTIAL,
        },
        "reviewer_2.core_field_review",
    )
    for reviewer_name, frames in (("reviewer_1", r1), ("reviewer_2", r2)):
        errors += validate_reviewer_values(
            frames["populated_field_review"],
            {"correct_yes_no_partial": YES_NO_PARTIAL, "supported_yes_no_partial": YES_NO_PARTIAL},
            f"{reviewer_name}.populated_field_review",
        )
        errors += validate_reviewer_values(
            frames["record_review"],
            {
                "source_recovered_yes_no": YES_NO,
                "sample_identity_correct_yes_no": YES_NO,
                "sample_count_correct_yes_no": YES_NO,
                "complete_core_row_correct_yes_no": YES_NO,
                "complete_populated_row_correct_yes_no": YES_NO,
                "recommended_audit_status": AUDIT_STATUSES,
            },
            f"{reviewer_name}.record_review",
        )
    if errors:
        raise ValueError("Reviewer workbook validation failed:\n" + "\n".join(errors))

    required_columns = {
        "core_field_review": [
            "present_in_source_yes_no",
            "extracted_value_correct_yes_no_partial",
        ],
        "populated_field_review": [
            "correct_yes_no_partial",
            "supported_yes_no_partial",
        ],
        "record_review": [
            "source_recovered_yes_no",
            "sample_identity_correct_yes_no",
            "sample_count_correct_yes_no",
            "complete_core_row_correct_yes_no",
            "complete_populated_row_correct_yes_no",
            "recommended_audit_status",
        ],
    }
    reviewer_missing: dict[str, dict[str, int]] = {}
    for reviewer_name, frames in (("reviewer_1", r1), ("reviewer_2", r2)):
        reviewer_missing[reviewer_name] = {
            f"{sheet}.{column}": int(frames[sheet][column].map(clean).eq("").sum())
            for sheet, columns in required_columns.items()
            for column in columns
        }
    total_reviewer_missing = sum(sum(values.values()) for values in reviewer_missing.values())

    key_map = {
        "core_field_review": ["validation_record_id", "field"],
        "populated_field_review": ["validation_record_id", "field"],
        "record_review": ["validation_record_id"],
    }
    input_map = {
        "core_field_review": [
            "present_in_source_yes_no",
            "gold_value",
            "extracted_value_correct_yes_no_partial",
            "evidence_supports_value_yes_no_partial",
            "notes",
        ],
        "populated_field_review": ["correct_yes_no_partial", "supported_yes_no_partial", "notes"],
        "record_review": [
            "source_recovered_yes_no",
            "sample_identity_correct_yes_no",
            "sample_count_correct_yes_no",
            "complete_core_row_correct_yes_no",
            "complete_populated_row_correct_yes_no",
            "recommended_audit_status",
            "notes",
        ],
    }
    tables: dict[str, pd.DataFrame] = {
        "instructions": pd.DataFrame(
            {
                "step": range(1, 5),
                "instruction": [
                    "Resolve only reviewer disagreements or missing required labels.",
                    "Use the source paper and frozen tolerance rules; do not inspect system status.",
                    "Complete every adjudicated field before scoring.",
                    "Do not alter validation_record_id or field columns.",
                ],
            }
        )
    }
    disagreement_counts: dict[str, int] = {}
    for sheet in sheets:
        keys = key_map[sheet]
        left = r1[sheet].copy()
        right = r2[sheet].copy()
        if left.duplicated(keys).any() or right.duplicated(keys).any():
            raise ValueError(f"Duplicate keys found in {sheet}")
        base_columns = [column for column in left.columns if column not in input_map[sheet]]
        left = left[base_columns + input_map[sheet]]
        right = right[keys + input_map[sheet]]
        left = left.rename(columns={column: f"r1_{column}" for column in input_map[sheet]})
        right = right.rename(columns={column: f"r2_{column}" for column in input_map[sheet]})
        merged = left.merge(right, on=keys, how="outer", validate="one_to_one", indicator=True)
        if not merged["_merge"].eq("both").all():
            raise ValueError(f"Reviewer row sets differ in {sheet}")
        merged = merged.drop(columns="_merge")
        for column in input_map[sheet]:
            merged[f"adjudicated_{column}"] = ""
        decision_columns = [column for column in input_map[sheet] if column != "notes" and column != "gold_value"]
        disagree = np.zeros(len(merged), dtype=bool)
        for column in decision_columns:
            disagree |= merged[f"r1_{column}"].map(label).ne(merged[f"r2_{column}"].map(label))
        merged.insert(len(base_columns), "reviewer_disagreement", np.where(disagree, "yes", "no"))
        disagreement_counts[sheet] = int(disagree.sum())
        tables[sheet] = merged
    tables["numeric_tolerances"] = tolerance_table()
    write_tables(out_path, tables, input_prefixes=("adjudicated_",))
    result = {
        "status": "awaiting_reviewer_completion" if total_reviewer_missing else "awaiting_adjudication",
        "reviewer_1": str(reviewer_1),
        "reviewer_2": str(reviewer_2),
        "adjudication_workbook": str(out_path),
        "disagreements": disagreement_counts,
        "required_reviewer_cells_missing": total_reviewer_missing,
        "missing_by_reviewer": reviewer_missing,
    }
    (out_path.parent / "merge_manifest.json").write_text(json.dumps(result, indent=2), encoding="utf-8")
    return result


def completion_report(adjudication: Path) -> dict[str, object]:
    core = read_sheet(adjudication, "core_field_review")
    populated = read_sheet(adjudication, "populated_field_review")
    record = read_sheet(adjudication, "record_review")
    requirements = {
        "core_field_review": {
            "adjudicated_present_in_source_yes_no": np.ones(len(core), dtype=bool),
            "adjudicated_gold_value": core["adjudicated_present_in_source_yes_no"].map(label).eq("yes").to_numpy(),
            "adjudicated_extracted_value_correct_yes_no_partial": np.ones(len(core), dtype=bool),
            "adjudicated_evidence_supports_value_yes_no_partial": core["extracted_value"].map(clean).ne("").to_numpy(),
        },
        "populated_field_review": {
            "adjudicated_correct_yes_no_partial": np.ones(len(populated), dtype=bool),
            "adjudicated_supported_yes_no_partial": np.ones(len(populated), dtype=bool),
        },
        "record_review": {
            "adjudicated_source_recovered_yes_no": np.ones(len(record), dtype=bool),
            "adjudicated_sample_identity_correct_yes_no": np.ones(len(record), dtype=bool),
            "adjudicated_sample_count_correct_yes_no": np.ones(len(record), dtype=bool),
            "adjudicated_complete_core_row_correct_yes_no": np.ones(len(record), dtype=bool),
            "adjudicated_complete_populated_row_correct_yes_no": np.ones(len(record), dtype=bool),
            "adjudicated_recommended_audit_status": np.ones(len(record), dtype=bool),
        },
    }
    frames = {"core_field_review": core, "populated_field_review": populated, "record_review": record}
    missing: dict[str, dict[str, int]] = {}
    for sheet, columns in requirements.items():
        frame = frames[sheet]
        sheet_missing: dict[str, int] = {}
        for column, required_mask in columns.items():
            if column not in frame:
                sheet_missing[column] = int(required_mask.sum())
            else:
                blank = frame[column].map(clean).eq("").to_numpy()
                sheet_missing[column] = int((blank & required_mask).sum())
        missing[sheet] = sheet_missing
    total_missing = sum(sum(items.values()) for items in missing.values())
    return {
        "status": "complete" if total_missing == 0 else "pending",
        "required_cells_missing": total_missing,
        "missing_by_sheet_and_column": missing,
    }


def safe_ratio(numerator: float, denominator: float) -> float | None:
    return float(numerator / denominator) if denominator else None


def primary_counts(frame: pd.DataFrame) -> tuple[int, int, int]:
    gold_present = frame["adjudicated_present_in_source_yes_no"].map(label).eq("yes")
    extracted_present = frame["extracted_value"].map(clean).ne("")
    correct = frame["adjudicated_extracted_value_correct_yes_no_partial"].map(label).eq("yes")
    true_positive = gold_present & extracted_present & correct
    false_positive = extracted_present & ~true_positive
    false_negative = gold_present & ~true_positive
    return int(true_positive.sum()), int(false_positive.sum()), int(false_negative.sum())


def primary_metrics(frame: pd.DataFrame) -> dict[str, float | int | None]:
    tp, fp, fn = primary_counts(frame)
    precision = safe_ratio(tp, tp + fp)
    recall = safe_ratio(tp, tp + fn)
    f1 = None
    if precision is not None and recall is not None and precision + recall:
        f1 = 2 * precision * recall / (precision + recall)
    return {"tp": tp, "fp": fp, "fn": fn, "precision": precision, "recall": recall, "f1": f1}


def cohen_kappa(values_1: pd.Series, values_2: pd.Series) -> tuple[float | None, float | None, int]:
    pairs = [(label(a), label(b)) for a, b in zip(values_1, values_2) if clean(a) and clean(b)]
    if not pairs:
        return None, None, 0
    n = len(pairs)
    agreement = sum(a == b for a, b in pairs) / n
    counts_1 = Counter(a for a, _ in pairs)
    counts_2 = Counter(b for _, b in pairs)
    categories = set(counts_1) | set(counts_2)
    expected = sum((counts_1[c] / n) * (counts_2[c] / n) for c in categories)
    kappa = (agreement - expected) / (1 - expected) if expected < 1 else 1.0
    return float(agreement), float(kappa), n


def normalize_numeric(field: str, value: float) -> float:
    if field == "JV_default_FF" and value > 1.5:
        return value / 100.0
    if field == "Stability_PCE_end_of_experiment" and value > 1.5:
        return value / 100.0
    return value


def numeric_agreement(core: pd.DataFrame) -> dict[str, object]:
    tolerance_fields = {field for field, _, _, _, _ in NUMERIC_TOLERANCES}
    assessed: list[bool] = []
    rows: list[dict[str, object]] = []
    for item in core.to_dict(orient="records"):
        field = clean(item.get("field"))
        if field not in tolerance_fields or label(item.get("adjudicated_present_in_source_yes_no")) != "yes":
            continue
        decision = label(item.get("adjudicated_extracted_value_correct_yes_no_partial"))
        if decision not in YES_NO_PARTIAL:
            continue
        passed = decision == "yes"
        assessed.append(passed)
        rows.append({"field": field, "passed": passed})
    by_field = []
    if rows:
        numeric_frame = pd.DataFrame(rows)
        for field, group in numeric_frame.groupby("field", sort=True):
            by_field.append({"field": field, "n": int(len(group)), "agreement": float(group["passed"].mean())})
    return {
        "n": len(assessed),
        "agreement": float(np.mean(assessed)) if assessed else None,
        "by_field": by_field,
        "decision_basis": "adjudicators applied frozen field tolerances after unit normalization",
    }


def categorical_macro_f1(core: pd.DataFrame) -> dict[str, object]:
    subset = core[
        core["field"].isin(CATEGORICAL_FIELDS)
        & core["adjudicated_present_in_source_yes_no"].map(label).eq("yes")
    ].copy()
    scores: list[float] = []
    class_rows: list[dict[str, object]] = []
    for field, group in subset.groupby("field", sort=True):
        gold = [canonical_category(field, value) for value in group["adjudicated_gold_value"]]
        predicted = [canonical_category(field, value) for value in group["extracted_value"]]
        for category in sorted(set(gold) | set(predicted)):
            tp = sum(g == category and p == category for g, p in zip(gold, predicted))
            fp = sum(g != category and p == category for g, p in zip(gold, predicted))
            fn = sum(g == category and p != category for g, p in zip(gold, predicted))
            precision = safe_ratio(tp, tp + fp) or 0.0
            recall = safe_ratio(tp, tp + fn) or 0.0
            score = 2 * precision * recall / (precision + recall) if precision + recall else 0.0
            scores.append(score)
            class_rows.append({"field": field, "category": category, "f1": score})
    return {"n": int(len(subset)), "macro_f1": float(np.mean(scores)) if scores else None, "classes": class_rows}


def summarize_group(frame: pd.DataFrame, group_column: str) -> pd.DataFrame:
    rows = []
    for group_name, group in frame.groupby(group_column, dropna=False, sort=True):
        metric = primary_metrics(group)
        rows.append({group_column: group_name, "n_fields": int(len(group)), **metric})
    return pd.DataFrame(rows)


def bootstrap_primary(core: pd.DataFrame, iterations: int, seed: int) -> dict[str, list[float] | None]:
    doi_map = core.groupby("Ref_DOI_number", dropna=False).indices
    groups = list(doi_map)
    rng = np.random.default_rng(seed)
    draws: dict[str, list[float]] = {"precision": [], "recall": [], "f1": []}
    for _ in range(iterations):
        sampled = rng.choice(groups, size=len(groups), replace=True)
        indices = np.concatenate([np.asarray(doi_map[group], dtype=int) for group in sampled])
        metric = primary_metrics(core.iloc[indices])
        for name in draws:
            if metric[name] is not None:
                draws[name].append(float(metric[name]))
    result: dict[str, list[float] | None] = {}
    for name, values in draws.items():
        result[name] = [float(x) for x in np.quantile(values, [0.025, 0.975])] if values else None
    return result


def add_textbox(slide, x, y, w, h, text_value, size=18, bold=False, color="20242A", align=PP_ALIGN.LEFT):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    frame = box.text_frame
    frame.clear()
    frame.word_wrap = True
    paragraph = frame.paragraphs[0]
    paragraph.alignment = align
    run = paragraph.add_run()
    run.text = text_value
    run.font.name = "Aptos"
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = RGBColor.from_string(color)
    return box


def build_results_pptx(metrics: dict[str, object], path: Path) -> None:
    presentation = Presentation()
    presentation.slide_width = Inches(13.333)
    presentation.slide_height = Inches(7.5)
    slide = presentation.slides.add_slide(presentation.slide_layouts[6])
    background = slide.background.fill
    background.solid()
    background.fore_color.rgb = RGBColor(255, 255, 255)
    add_textbox(slide, 0.55, 0.3, 12.2, 0.55, "Structured-extraction validation", 26, True)
    add_textbox(slide, 0.55, 0.88, 12.2, 0.35, "Frozen 60-record benchmark; 95% intervals use DOI-cluster bootstrap", 12, False, "5D646F")
    primary = metrics["primary_field_metrics"]
    record = metrics["record_metrics"]
    panels = [
        ("Field-level value validity", [("Precision", primary["precision"]), ("Recall", primary["recall"]), ("F1", primary["f1"]), ("Evidence support", metrics["evidence_support_validity"])]),
        ("Complete-record validity", [("Sample linking", record["sample_linking_accuracy"]), ("Core row", record["complete_core_row_accuracy"]), ("Populated row", record["complete_populated_row_accuracy"]), ("Audit decision", record["audit_decision_accuracy"])]),
    ]
    for panel_index, (title, values) in enumerate(panels):
        x0 = 0.65 + panel_index * 6.35
        shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x0), Inches(1.45), Inches(5.95), Inches(4.85))
        shape.fill.solid()
        shape.fill.fore_color.rgb = RGBColor.from_string("F5F5F3")
        shape.line.color.rgb = RGBColor.from_string("20242A")
        add_textbox(slide, x0 + 0.25, 1.68, 5.45, 0.4, title, 17, True)
        for row_index, (name, value) in enumerate(values):
            y = 2.3 + row_index * 0.88
            numeric_value = float(value) if value is not None else 0.0
            add_textbox(slide, x0 + 0.25, y, 1.65, 0.32, name, 12, False)
            track = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x0 + 1.95), Inches(y + 0.03), Inches(2.9), Inches(0.25))
            track.fill.solid()
            track.fill.fore_color.rgb = RGBColor.from_string("D9D9D6")
            track.line.fill.background()
            bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x0 + 1.95), Inches(y + 0.03), Inches(2.9 * numeric_value), Inches(0.25))
            bar.fill.solid()
            bar.fill.fore_color.rgb = RGBColor.from_string("2A9D8F")
            bar.line.fill.background()
            add_textbox(slide, x0 + 5.0, y - 0.04, 0.65, 0.35, f"{numeric_value:.2f}", 13, True, align=PP_ALIGN.RIGHT)
    add_textbox(slide, 0.65, 6.58, 12.0, 0.42, "Primary scoring counts partial decisions as incorrect; secondary partial-credit results remain in the data table.", 11, False, "5D646F")
    path.parent.mkdir(parents=True, exist_ok=True)
    presentation.save(path)


def replace_paragraph_text(paragraph, text_value: str) -> None:
    for run in paragraph.runs:
        run.text = ""
    run = paragraph.runs[0] if paragraph.runs else paragraph.add_run()
    run.text = text_value


def integrate_manuscript(template: Path, output: Path, metrics: dict[str, object]) -> None:
    document = Document(template)
    primary = metrics["primary_field_metrics"]
    ci = metrics["doi_cluster_bootstrap_95ci"]
    record = metrics["record_metrics"]
    evidence = metrics["evidence_support_validity"]
    numeric = metrics["numeric_agreement"]["agreement"]
    agreement = metrics["reviewer_agreement"]["pooled_decisions"]
    result_text = (
        f"The blinded 60-record structured-extraction benchmark produced field-value precision of {primary['precision']:.3f} "
        f"(95% DOI-cluster bootstrap interval {ci['precision'][0]:.3f}-{ci['precision'][1]:.3f}), recall of {primary['recall']:.3f} "
        f"({ci['recall'][0]:.3f}-{ci['recall'][1]:.3f}), and F1 of {primary['f1']:.3f} ({ci['f1'][0]:.3f}-{ci['f1'][1]:.3f}). "
        f"Evidence-support validity was {evidence:.3f}, numeric tolerance agreement was {numeric:.3f}, sample-linking accuracy was "
        f"{record['sample_linking_accuracy']:.3f}, and complete core-row accuracy was {record['complete_core_row_accuracy']:.3f}. "
        f"Reviewer agreement across pooled field decisions was {agreement['percent_agreement']:.3f} with Cohen kappa {agreement['cohen_kappa']:.3f}."
    )
    for paragraph in document.paragraphs:
        if "[RESULT REQUIRED BEFORE SUBMISSION]" in paragraph.text:
            replace_paragraph_text(paragraph, result_text)
        elif "A frozen, two-reviewer 60-record structured-extraction benchmark has been prepared" in paragraph.text:
            sentence = (
                f"In a frozen two-reviewer 60-record benchmark, field-value precision, recall, and F1 were "
                f"{primary['precision']:.3f}, {primary['recall']:.3f}, and {primary['f1']:.3f}, respectively."
            )
            replace_paragraph_text(
                paragraph,
                paragraph.text.replace(
                    "A frozen, two-reviewer 60-record structured-extraction benchmark has been prepared, but its adjudicated field-level results must be inserted before submission.",
                    sentence,
                ),
            )
        elif "Final claims about extraction quality remain contingent on the blinded 60-record human benchmark." in paragraph.text:
            replace_paragraph_text(
                paragraph,
                paragraph.text.replace(
                    "Final claims about extraction quality remain contingent on the blinded 60-record human benchmark.",
                    f"The blinded benchmark yielded a conservative field-value F1 of {primary['f1']:.3f}, with complete-row performance reported separately.",
                ),
            )
    metric_values = {
        "Field-value precision": primary["precision"],
        "Field-value recall": primary["recall"],
        "Field-value F1": primary["f1"],
        "Evidence-support validity": evidence,
        "Numeric tolerance agreement": numeric,
        "Categorical macro-F1": metrics["categorical_macro_f1"]["macro_f1"],
        "Sample-linking accuracy": record["sample_linking_accuracy"],
        "Complete core-row accuracy": record["complete_core_row_accuracy"],
        "Complete populated-row accuracy": record["complete_populated_row_accuracy"],
        "Audit-decision agreement": record["audit_decision_accuracy"],
    }
    for table in document.tables:
        if table.rows and clean(table.cell(0, 0).text) == "Metric" and any("Pending blinded adjudication" in cell.text for row in table.rows for cell in row.cells):
            for row in table.rows[1:]:
                name = clean(row.cells[0].text)
                if name in metric_values and metric_values[name] is not None:
                    row.cells[1].text = f"{float(metric_values[name]):.3f}"
    output.parent.mkdir(parents=True, exist_ok=True)
    document.save(output)


def score(adjudication: Path, master: Path, out_dir: Path, iterations: int, manuscript: Path | None) -> dict[str, object]:
    out_dir.mkdir(parents=True, exist_ok=True)
    completion = completion_report(adjudication)
    (out_dir / "annotation_completion_report.json").write_text(json.dumps(completion, indent=2), encoding="utf-8")
    if completion["status"] != "complete":
        return completion

    core = read_sheet(adjudication, "core_field_review")
    populated = read_sheet(adjudication, "populated_field_review")
    record = read_sheet(adjudication, "record_review")
    master_manifest = read_sheet(master, "sample_manifest")
    system_key = master_manifest[
        ["validation_record_id", "validation_stratum", "integration_status", "Ref_DOI_number"]
    ].drop_duplicates("validation_record_id")
    core = core.drop(columns=[column for column in ("Ref_DOI_number",) if column in core]).merge(system_key, on="validation_record_id", how="left", validate="many_to_one")
    populated = populated.drop(columns=[column for column in ("Ref_DOI_number",) if column in populated]).merge(system_key, on="validation_record_id", how="left", validate="many_to_one")
    record = record.drop(columns=[column for column in ("Ref_DOI_number",) if column in record]).merge(system_key, on="validation_record_id", how="left", validate="one_to_one")

    primary = primary_metrics(core)
    bootstrap = bootstrap_primary(core, iterations=iterations, seed=42)
    core_assessed = (
        core["extracted_value"].map(clean).ne("")
        & core["adjudicated_evidence_supports_value_yes_no_partial"].map(label).isin(YES_NO_PARTIAL)
    )
    populated_assessed = populated["adjudicated_supported_yes_no_partial"].map(label).isin(YES_NO_PARTIAL)
    supported_yes = int(
        core.loc[core_assessed, "adjudicated_evidence_supports_value_yes_no_partial"].map(label).eq("yes").sum()
        + populated.loc[populated_assessed, "adjudicated_supported_yes_no_partial"].map(label).eq("yes").sum()
    )
    unsupported_no = int(
        core.loc[core_assessed, "adjudicated_evidence_supports_value_yes_no_partial"].map(label).eq("no").sum()
        + populated.loc[populated_assessed, "adjudicated_supported_yes_no_partial"].map(label).eq("no").sum()
    )
    evidence_assessed_n = int(core_assessed.sum() + populated_assessed.sum())
    evidence_support = safe_ratio(supported_yes, evidence_assessed_n)
    unsupported_rate = safe_ratio(unsupported_no, evidence_assessed_n)
    numeric = numeric_agreement(core)
    categorical = categorical_macro_f1(core)

    system_status = record["integration_status"].map(label).replace({"": "accepted"})
    gold_status = record["adjudicated_recommended_audit_status"].map(label)
    record_metrics = {
        "n_records": int(len(record)),
        "sample_linking_accuracy": float(record["adjudicated_sample_identity_correct_yes_no"].map(label).eq("yes").mean()),
        "sample_count_accuracy": float(record["adjudicated_sample_count_correct_yes_no"].map(label).eq("yes").mean()),
        "complete_core_row_accuracy": float(record["adjudicated_complete_core_row_correct_yes_no"].map(label).eq("yes").mean()),
        "complete_populated_row_accuracy": float(record["adjudicated_complete_populated_row_correct_yes_no"].map(label).eq("yes").mean()),
        "audit_decision_accuracy": float(system_status.eq(gold_status).mean()),
    }

    agreement_rows: list[dict[str, object]] = []
    agreement_specs = [
        ("core_present", core, "r1_present_in_source_yes_no", "r2_present_in_source_yes_no"),
        ("core_correct", core, "r1_extracted_value_correct_yes_no_partial", "r2_extracted_value_correct_yes_no_partial"),
        ("core_evidence", core, "r1_evidence_supports_value_yes_no_partial", "r2_evidence_supports_value_yes_no_partial"),
        ("populated_correct", populated, "r1_correct_yes_no_partial", "r2_correct_yes_no_partial"),
        ("populated_evidence", populated, "r1_supported_yes_no_partial", "r2_supported_yes_no_partial"),
        ("sample_identity", record, "r1_sample_identity_correct_yes_no", "r2_sample_identity_correct_yes_no"),
        ("audit_status", record, "r1_recommended_audit_status", "r2_recommended_audit_status"),
    ]
    pooled_1: list[object] = []
    pooled_2: list[object] = []
    for name, frame, first, second in agreement_specs:
        percent, kappa, n = cohen_kappa(frame[first], frame[second])
        agreement_rows.append({"decision": name, "n": n, "percent_agreement": percent, "cohen_kappa": kappa})
        for value_1, value_2 in zip(frame[first], frame[second]):
            if clean(value_1) and clean(value_2):
                pooled_1.append(value_1)
                pooled_2.append(value_2)
    pooled_percent, pooled_kappa, pooled_n = cohen_kappa(pd.Series(pooled_1), pd.Series(pooled_2))
    agreement_summary = {
        "pooled_decisions": {
            "n": pooled_n,
            "percent_agreement": pooled_percent,
            "cohen_kappa": pooled_kappa,
        },
        "by_decision": agreement_rows,
    }

    metrics = {
        "status": "complete",
        "benchmark_records": int(len(record)),
        "benchmark_dois": int(record["Ref_DOI_number"].nunique()),
        "primary_field_metrics": primary,
        "doi_cluster_bootstrap_95ci": bootstrap,
        "evidence_support_validity": evidence_support,
        "unsupported_value_rate": unsupported_rate,
        "numeric_agreement": numeric,
        "categorical_macro_f1": categorical,
        "record_metrics": record_metrics,
        "reviewer_agreement": agreement_summary,
        "partial_policy": "partial counted as incorrect in primary metrics",
    }
    (out_dir / "structured_extraction_metrics.json").write_text(json.dumps(metrics, indent=2), encoding="utf-8")
    summary_rows = []
    for metric_name in ("precision", "recall", "f1"):
        interval = bootstrap[metric_name]
        summary_rows.append(
            {
                "metric": f"field_value_{metric_name}",
                "value": primary[metric_name],
                "ci_low": interval[0] if interval else None,
                "ci_high": interval[1] if interval else None,
                "n": len(core),
            }
        )
    for metric_name, value in (
        ("evidence_support_validity", evidence_support),
        ("unsupported_value_rate", unsupported_rate),
        ("numeric_tolerance_agreement", numeric["agreement"]),
        ("categorical_macro_f1", categorical["macro_f1"]),
        *record_metrics.items(),
    ):
        if metric_name != "n_records":
            summary_rows.append({"metric": metric_name, "value": value, "ci_low": None, "ci_high": None, "n": None})
    summary = pd.DataFrame(summary_rows)
    by_family = summarize_group(core, "field_family")
    by_stratum = summarize_group(core, "validation_stratum")
    agreement_frame = pd.DataFrame(agreement_rows)
    summary.to_csv(out_dir / "structured_extraction_metrics.csv", index=False)
    by_family.to_csv(out_dir / "metrics_by_field_family.csv", index=False)
    by_stratum.to_csv(out_dir / "metrics_by_validation_stratum.csv", index=False)
    agreement_frame.to_csv(out_dir / "reviewer_agreement.csv", index=False)
    with pd.ExcelWriter(out_dir / "structured_extraction_validation_results.xlsx", engine="openpyxl") as writer:
        summary.to_excel(writer, sheet_name="summary", index=False)
        by_family.to_excel(writer, sheet_name="by_field_family", index=False)
        by_stratum.to_excel(writer, sheet_name="by_stratum", index=False)
        agreement_frame.to_excel(writer, sheet_name="reviewer_agreement", index=False)
        pd.DataFrame(numeric["by_field"]).to_excel(writer, sheet_name="numeric_by_field", index=False)
    snippet = (
        f"The frozen 60-record benchmark yielded conservative field-value precision of {primary['precision']:.3f} "
        f"(95% DOI-cluster bootstrap interval {bootstrap['precision'][0]:.3f}-{bootstrap['precision'][1]:.3f}), "
        f"recall of {primary['recall']:.3f} ({bootstrap['recall'][0]:.3f}-{bootstrap['recall'][1]:.3f}), and F1 of "
        f"{primary['f1']:.3f} ({bootstrap['f1'][0]:.3f}-{bootstrap['f1'][1]:.3f}). Evidence-support validity was "
        f"{evidence_support:.3f}; sample-linking and complete core-row accuracy were {record_metrics['sample_linking_accuracy']:.3f} "
        f"and {record_metrics['complete_core_row_accuracy']:.3f}. Partial decisions were counted as incorrect in the primary analysis."
    )
    (out_dir / "manuscript_results_snippet.md").write_text(snippet + "\n", encoding="utf-8")
    build_results_pptx(metrics, out_dir / "Figure_9_structured_extraction_validation_editable.pptx")
    if manuscript:
        integrated = out_dir / f"{manuscript.stem}_with_structured_validation.docx"
        integrate_manuscript(manuscript, integrated, metrics)
        metrics["integrated_manuscript"] = str(integrated)
    return metrics


def main() -> None:
    parser = argparse.ArgumentParser(description="Build, merge, and score the frozen structured-extraction benchmark.")
    subparsers = parser.add_subparsers(dest="command", required=True)

    build_parser = subparsers.add_parser("build", help="Create blinded reviewer packages from the frozen master workbook.")
    build_parser.add_argument("--master", type=Path, default=DEFAULT_MASTER)
    build_parser.add_argument("--out-dir", type=Path, default=DEFAULT_OUT)

    merge_parser = subparsers.add_parser("merge", help="Merge completed blinded reviews into an adjudication workbook.")
    merge_parser.add_argument("--reviewer-1", type=Path, required=True)
    merge_parser.add_argument("--reviewer-2", type=Path, required=True)
    merge_parser.add_argument("--out", type=Path, default=DEFAULT_OUT / "structured_extraction_adjudication.xlsx")

    score_parser = subparsers.add_parser("score", help="Score a completed adjudication workbook.")
    score_parser.add_argument("--adjudication", type=Path, required=True)
    score_parser.add_argument("--master", type=Path, default=DEFAULT_MASTER)
    score_parser.add_argument("--out-dir", type=Path, default=DEFAULT_OUT / "results")
    score_parser.add_argument("--bootstrap-iterations", type=int, default=2000)
    score_parser.add_argument("--manuscript", type=Path)

    args = parser.parse_args()
    if args.command == "build":
        result = build_blinded(args.master.resolve(), args.out_dir.resolve())
    elif args.command == "merge":
        result = merge_reviews(args.reviewer_1.resolve(), args.reviewer_2.resolve(), args.out.resolve())
    else:
        result = score(
            args.adjudication.resolve(),
            args.master.resolve(),
            args.out_dir.resolve(),
            args.bootstrap_iterations,
            args.manuscript.resolve() if args.manuscript else None,
        )
    print(json.dumps(result, indent=2))


if __name__ == "__main__":
    main()
